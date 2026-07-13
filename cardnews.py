# -*- coding: utf-8 -*-
"""
IamI 여자 골격진단 인스타 카드뉴스 생성기 (세로형 4:5, 9장)
사용법: python3 cardnews.py template.pptx spec.json photos_dir out_dir
photos_dir 에는 01.png(커버), 02.png ~ 08.png (본문 슬라이드 2~8) 필요 (png/jpg 모두 가능)
"""
import sys, os, json, glob, copy, zipfile, subprocess, shutil
from pptx import Presentation
from pptx.util import Emu, Inches
from PIL import Image, ImageFont

EMU = 914400
PINK = (236, 95, 184)
F5 = '/root/.fonts/Paperlogy-5Medium.ttf'
F7 = '/root/.fonts/Paperlogy-7Bold.ttf'

def text_w_in(text, font_path, pt):
    f = ImageFont.truetype(font_path, int(round(pt * 96 / 72)))
    return f.getlength(text) / 96.0

def parse_segments(s):
    """'a [b] c' -> [('a ',False),('[b]',True),(' c',False)]"""
    out, buf, i = [], '', 0
    while i < len(s):
        if s[i] == '[':
            j = s.find(']', i)
            if j == -1:
                buf += s[i:]; break
            if buf: out.append((buf, False)); buf = ''
            out.append((s[i:j+1], True)); i = j + 1
        else:
            buf += s[i]; i += 1
    if buf: out.append((buf, False))
    return out

def set_runs(para, segments, dark_rpr, pink_rpr):
    """문단의 런을 세그먼트대로 재구성 (서식은 rPr 클론)"""
    from lxml import etree
    p = para._p
    for r in list(p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r')):
        p.remove(r)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for text, is_pink in segments:
        r = etree.SubElement(p, f'{{{ns}}}r')
        rpr = copy.deepcopy(pink_rpr if is_pink else dark_rpr)
        r.append(rpr)
        t = etree.SubElement(r, f'{{{ns}}}t')
        t.text = text
        t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

def get_rprs(para):
    """문단에서 (dark_rpr, pink_rpr) 추출; 핑크 없으면 dark 복제 후 색 변경"""
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    dark = pink = None
    for r in para.runs:
        rpr = r._r.find(f'{{{ns}}}rPr')
        if rpr is None: continue
        c = r.font.color
        is_pink = c and c.type is not None and tuple(c.rgb) == PINK
        if is_pink and pink is None: pink = copy.deepcopy(rpr)
        if not is_pink and dark is None: dark = copy.deepcopy(rpr)
    if dark is None: dark = copy.deepcopy(pink)
    if pink is None:
        pink = copy.deepcopy(dark)
        fill = pink.find(f'{{{ns}}}solidFill')
        if fill is not None: pink.remove(fill)
        sf = etree.SubElement(pink, f'{{{ns}}}solidFill')
        clr = etree.SubElement(sf, f'{{{ns}}}srgbClr')
        clr.set('val', 'EC5FB8')
        pink.insert(0, sf)  # solidFill must come early; reinsert properly
        pink.remove(sf); 
        # correct order: insert after ln if any, keep simple: append works for LO/PPT tolerance
        pink.append(sf)
    return dark, pink

def set_text_single(shape, text):
    """단일 런 텍스트박스: 첫 런에 텍스트만 교체"""
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        for r in para.runs:
            r.text = text if first else ''
            first = False

def find_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name: return sh
    return None

def shapes_sorted(slide):
    return list(slide.shapes)

def main():
    tpl, specf, photos, outdir = sys.argv[1:5]
    spec = json.load(open(specf, encoding='utf-8'))
    os.makedirs(outdir, exist_ok=True)
    work = os.path.join(outdir, '_work'); shutil.rmtree(work, ignore_errors=True); os.makedirs(work)

    # 1) 사진을 PNG로 정규화하여 미디어 교체 준비
    photo_map = {}  # media filename -> new png path
    slots = {1: 'image1.png', 2: 'image3.png', 3: 'image5.png', 4: 'image6.png',
             5: 'image7.png', 6: 'image8.png', 7: 'image9.png', 8: 'image10.png'}
    ASPECT = {1: 11.25 / 14.0625}  # 커버 4:5, 본문은 9.92x5.83
    for idx, media in slots.items():
        cand = sorted(glob.glob(os.path.join(photos, f'{idx:02d}*')))
        if not cand:
            print(f'경고: 사진 {idx:02d} 없음 — 기존 이미지 유지'); continue
        im = Image.open(cand[0]).convert('RGB')
        target = ASPECT.get(idx, 9.92 / 5.83)
        w, h = im.size
        cur = w / h
        if abs(cur - target) > 0.01:  # 센터 크롭
            if cur > target:
                nw = int(h * target); x = (w - nw) // 2
                im = im.crop((x, 0, x + nw, h))
            else:
                nh = int(w / target); y = (h - nh) // 2
                im = im.crop((0, y, w, y + nh))
        if idx == 1:
            # 커버: 은은한 다크 그라데이션 (배경이 비쳐 보이는 수준, 텍스트 가독성용)
            # 원본 사진에 텍스트가 박혀있으면 안 됨 — 반드시 클린 원본 사용
            cw, ch = im.size
            ov = Image.new('L', (1, ch), 0)
            opx = ov.load()
            for yy in range(ch):
                f = yy / ch
                if f < 0.38: a = 0
                elif f < 0.60: a = int(115 * (f - 0.38) / 0.22)
                else: a = 115  # 최대 45% 어둡기 — 배경 비침 유지
                opx[0, yy] = a
            ov = ov.resize((cw, ch))
            im = Image.composite(Image.new('RGB', (cw, ch), (10, 10, 14)), im, ov)
        p = os.path.join(work, media)
        im.save(p, 'PNG')
        photo_map[media] = p

    prs = Presentation(tpl)
    S = list(prs.slides)

    # ---------- 슬라이드 1 (커버) ----------
    s = S[0]
    cov = spec['cover']
    sh = shapes_sorted(s)
    # 상단 우측 섹션명 (TextBox 5, pos x~5.94)
    for shp in sh:
        if shp.has_text_frame and abs(Emu(shp.left).inches - 5.94) < 0.05 and Emu(shp.top).inches < 1:
            set_text_single(shp, cov['section'])
    # 라벨 필 텍스트 (TextBox 7)
    lb = find_by_name(s, 'TextBox 7')
    if lb and cov.get('label'): set_text_single(lb, cov['label'])
    # 타이틀 3줄 (TextBox 8,9,10)
    title_boxes = [find_by_name(s, 'TextBox 8'), find_by_name(s, 'TextBox 9'), find_by_name(s, 'TextBox 10')]
    lines = cov['title_lines']
    for box, line in zip(title_boxes, lines):
        para = box.text_frame.paragraphs[0]
        dark, pink = get_rprs(para)
        set_runs(para, parse_segments(line), dark, pink)
    # 밑줄 바: 마지막 줄의 핑크 세그먼트 아래
    bar = find_by_name(s, 'Rounded Rectangle 11')
    segs = parse_segments(lines[-1])
    x = 0.67; w = None; acc = 0.0
    for text, is_pink in segs:
        tw = text_w_in(text, F7, 45)
        if is_pink and w is None:
            x = 0.67 + acc; w = tw
        acc += tw
    if w:
        bar.left = Inches(x); bar.width = Inches(w)
    bar.top = Inches(10.35)  # 1차 렌더 측정 존(8.83~9.79in) 밖에 주차 — 2차에서 실측 재배치
    # 해시태그 필 (최대 10개, 2열 자동 배치)
    tags = cov['hashtags'][:10]
    pill_rects, pill_texts = [], []
    for shp in sh:
        t = Emu(shp.top).inches
        if 10.9 < t < 12.2 and Emu(shp.height).inches < 0.6:
            if shp.has_text_frame and shp.shape_type == 17:
                pill_texts.append(shp)
            elif shp.shape_type == 1:
                pill_rects.append(shp)
    pill_rects.sort(key=lambda p: (Emu(p.top).inches, Emu(p.left).inches))
    pill_texts.sort(key=lambda p: (Emu(p.top).inches, Emu(p.left).inches))
    GAP, X0, XMAX = 0.14, 0.67, 10.59
    Y1, Y2 = 10.96, 11.58
    x_cur, y_cur = X0, Y1
    for i in range(len(pill_rects)):
        if i < len(tags):
            w = text_w_in(tags[i], F5, 17.25) + 0.378
            if x_cur + w > XMAX and y_cur == Y1:
                x_cur, y_cur = X0, Y2
            for obj in (pill_rects[i], pill_texts[i]):
                obj.left = Inches(x_cur); obj.top = Inches(y_cur); obj.width = Inches(w)
            set_text_single(pill_texts[i], tags[i])
            x_cur += w + GAP
        else:
            for obj in (pill_rects[i], pill_texts[i]):
                sp = obj._element; sp.getparent().remove(sp)
    # 하단 캡션 (TextBox 32)
    cap = find_by_name(s, 'TextBox 32')
    if cap: set_text_single(cap, cov['caption'])

    # ---------- 슬라이드 2~8 (본문) ----------
    for si, sl_spec in enumerate(spec['slides'], start=1):
        s = S[si]
        # 섹션명
        sec = find_by_name(s, 'TextBox 4')
        set_text_single(sec, sl_spec['section'])
        # 타이틀 (TextBox 5)
        tb = find_by_name(s, 'TextBox 5')
        para = tb.text_frame.paragraphs[0]
        dark, pink = get_rprs(para)
        set_runs(para, parse_segments(sl_spec['title']), dark, pink)
        # 밑줄 바를 핑크 [괄호] 세그먼트 바로 아래로 (35.25pt Paperlogy 7 Bold)
        bar2 = find_by_name(s, 'Rounded Rectangle 6')
        if bar2 is not None:
            acc = 0.0; bx = None; bw = None
            for text, is_pink in parse_segments(sl_spec['title']):
                tw = text_w_in(text, F7, 35.25)
                if is_pink and bw is None:
                    bx = 0.67 + acc; bw = tw
                acc += tw
            if bw:
                bar2.left = Inches(bx); bar2.width = Inches(bw)
        # 서브타이틀 (TextBox 7)
        set_text_single(find_by_name(s, 'TextBox 7'), sl_spec['subtitle'])
        # 불릿 3개 (TextBox 10, 12, 14)
        for name, txt in zip(['TextBox 10', 'TextBox 12', 'TextBox 14'], sl_spec['bullets']):
            set_text_single(find_by_name(s, name), txt)
        # 핵심 한 줄 (TextBox 18)
        set_text_single(find_by_name(s, 'TextBox 18'), sl_spec['key_line'])
        # 푸터 캡션 (TextBox 19)
        set_text_single(find_by_name(s, 'TextBox 19'), sl_spec['caption'])

    # ---------- 슬라이드 9 (아웃트로) ----------
    s = S[8]
    out = spec.get('outro', {})
    if out.get('section'):
        set_text_single(find_by_name(s, 'TextBox 5'), out['section'])
    if out.get('subtitle'):
        set_text_single(find_by_name(s, 'TextBox 8'), out['subtitle'])

    mid = os.path.join(work, 'mid.pptx')
    prs.save(mid)

    # 2) 미디어 교체 (zip 재작성)
    final_pptx = os.path.join(work, 'final.pptx')
    with zipfile.ZipFile(mid) as zin, zipfile.ZipFile(final_pptx, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            base = os.path.basename(item)
            if item.startswith('ppt/media/') and base in photo_map:
                zout.write(photo_map[base], item)
            else:
                zout.writestr(item, zin.read(item))

    # 3) 1차 렌더링 → 핑크 텍스트 실측 → 밑줄 바 재배치 → 2차 렌더링
    import numpy as np
    subprocess.run(['python', '/mnt/skills/public/pptx/scripts/office/soffice.py',
                    '--headless', '--convert-to', 'pdf', '--outdir', work, final_pptx],
                   check=True, capture_output=True)
    pdf = os.path.join(work, 'final.pdf')
    subprocess.run(['pdftoppm', '-jpeg', '-r', '96', pdf, os.path.join(work, 'pass1')], check=True)

    def pink_bbox(img_path, y0, y1):
        arr = np.array(Image.open(img_path).convert('RGB')).astype(int)
        d = np.abs(arr - np.array([236, 95, 184])).sum(axis=2)
        m = d < 60
        m[:y0, :] = False; m[y1:, :] = False
        ys, xs = np.nonzero(m)
        if len(xs) == 0: return None
        return xs.min(), xs.max(), ys.min(), ys.max()

    GAP_PX = 7
    prs2 = Presentation(final_pptx)
    S2 = list(prs2.slides)
    # 커버: 마지막 타이틀 줄(8.83~9.66in) 핑크 실측
    bb = pink_bbox(os.path.join(work, 'pass1-1.jpg'), 848, 940)
    bar = find_by_name(S2[0], 'Rounded Rectangle 11')
    if bb and bar is not None:
        x0, x1, ty0, ty1 = bb
        bar.left = Emu(int(x0 / 96 * EMU)); bar.width = Emu(int((x1 - x0 + 1) / 96 * EMU))
        bar.top = Emu(int((ty1 + GAP_PX) / 96 * EMU))
    # 본문 2~8: 타이틀 존(y 120~200px) 핑크 실측
    for si in range(1, 8):
        bb = pink_bbox(os.path.join(work, f'pass1-{si+1}.jpg'), 118, 202)
        bar = find_by_name(S2[si], 'Rounded Rectangle 6')
        if bb and bar is not None:
            x0, x1, ty0, ty1 = bb
            bar.left = Emu(int(x0 / 96 * EMU)); bar.width = Emu(int((x1 - x0 + 1) / 96 * EMU))
            bar.top = Emu(int((ty1 + GAP_PX) / 96 * EMU))
    final2 = os.path.join(work, 'final2.pptx')
    prs2.save(final2)

    subprocess.run(['python', '/mnt/skills/public/pptx/scripts/office/soffice.py',
                    '--headless', '--convert-to', 'pdf', '--outdir', work, final2],
                   check=True, capture_output=True)
    pdf = os.path.join(work, 'final2.pdf')
    subprocess.run(['pdftoppm', '-jpeg', '-r', '96', pdf, os.path.join(work, 'card')], check=True)
    prefix = spec.get('file_prefix', '카드')
    names = spec.get('file_names', [])
    outs = []
    for i in range(1, 10):
        src = os.path.join(work, f'card-{i}.jpg')
        im = Image.open(src).convert('RGB').resize((1080, 1350), Image.LANCZOS)
        label = names[i-1] if i-1 < len(names) else f'{i:02d}'
        dst = os.path.join(outdir, f'{i:02d}_{label}.jpg' if not label.startswith(f'{i:02d}') else f'{label}.jpg')
        im.save(dst, 'JPEG', quality=92)
        outs.append(dst)
    print('\n'.join(outs))

if __name__ == '__main__':
    main()
